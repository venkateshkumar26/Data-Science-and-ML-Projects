from pptx import Presentation
from typing import List, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt(outline:Dict[str,Any],image_paths:List[str],put_path:str="generated.pptx") -> Presentation:
    prs = Presentation()
    slides=outline.get("slides",[])
    if not slides:
        raise ValueError("Invalid outline: missing 'slides' array")
    title_slide=prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text=slides[0].get("title","Untitled")
    if len(title_slide.placeholders) > 1:
        try:
            title_slide.placeholders[1].text="Auto Generated Slides"
        except:
            pass
    for idx,slide_data in enumerate(slides[1:]):
        layout=prs.slide_layouts[1]
        slide=prs.slides.add_slide(layout)
        slide.shapes.title.text=slide_data.get("title","")
        bullets:List[str,Any]=slide_data.get("bullets",[])
        content=slide.placeholders[1].text_frame
        content.clear()
        if bullets:
            first=content.paragraphs[0]
            first.text=bullets[0]
            first.level=0
            for b in bullets[1:]:
                p=content.add_paragraph()
                p.text=b
                p.level=1

        image_path=image_paths.get(idx,"")
        if image_path:
            try:
                slide.shapes.add_picture(image_path,left=Inches(1),top=Inches(1),width=Inches(3.3))
            except Exception as e:
                print(f"[WARN] Failed to add image to slide {idx}: {e}")
        prs.save(put_path)
    return put_path


    